# Demo path = "/Users/PinskyRobin/Desktop/demo.pptx"
from pptx import Presentation
import sys


def ObtainPptNote(ppt_path):
    cnt = 1  # 标记页数
    obj = Presentation(ppt_path)  # 得到路径下的ppt文件对象
    fp = open('notes.txt', 'wt')
    saved_stdout = sys.stdout  # 将系统标准输出流保存
    sys.stdout = open('notes.txt', 'wt')  # 将标准输出流重定向,将输出流指向该txt文件
    notes = []  # 定义存储备注的列表
    for index, slide in enumerate(obj.slides, start=1):
        text = slide.notes_slide.notes_text_frame.text  # 提取备注
        if text != '':  # 备注非空
            print(text)  # 输出备注至txt文件
            notes.append(text)  # 将备注添加至列表中
        else:
            print('**No notes in page', cnt)  # 提示此页无备注
            notes.append('\n')
        cnt = cnt + 1  # 计数器自增
    sys.stdout = saved_stdout  # 还原标准输出流
    return notes  # 返回一个存储备注的列表


if __name__ == '__main__':
    print("请输入ppt所在的文件目录：")
    ppt_path = input()  # 获取pptx文件的存储路径
    print(ObtainPptNote(ppt_path))  # 获取备注
