from pptx import Presentation


def ObtainPptNote(ppt_path):
    obj = Presentation(ppt_path)  # 得到路径下的ppt文件对象
    notes = []  # 定义存储备注的列表
    for index, slide in enumerate(obj.slides, start=1):
        text = slide.notes_slide.notes_text_frame.text  # 提取备注
        if text != '':  # 备注非空
            notes.append(text)  # 将备注添加至列表中
        else:
            notes.append('\n')
    return notes  # 返回一个存储备注的列表

