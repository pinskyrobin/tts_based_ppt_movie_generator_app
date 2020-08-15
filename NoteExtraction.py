from pptx import Presentation

# 返回包含pptx文件备注的列表
def ObtainPptNote(ppt_path):
    if ppt_path.strip() == '':  # 空路径,直接返回
        return []
    obj = Presentation(ppt_path)  # 得到路径下的PPT文件对象
    notes = []
    for index, slide in enumerate(obj.slides, start=1):
        text = slide.notes_slide.notes_text_frame.text  # 提取备注
        notes.append(text)  # 将备注添加至列表中
    return notes
